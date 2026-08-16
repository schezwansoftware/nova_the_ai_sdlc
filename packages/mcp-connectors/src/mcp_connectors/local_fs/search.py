"""Shared local-filesystem directory-walk / text-search / path-safety
logic for the `local_docs` and `onedrive` connectors. Both are pure
local-filesystem access -- no HTTP client, no OAuth, no credential of any
kind (`onedrive` reads the OneDrive desktop client's already-synced local
folder directly, rather than calling Microsoft Graph -- see
`mcp_connectors/onedrive/client.py`'s module docstring for why that was a
deliberate choice, not a shortcut). This module holds the logic genuinely
shared between the two; each connector's own `client.py` is a thin
wrapper that supplies its allowlisted directories, its `source` label,
its configured `file_categories`, and whether OneDrive-specific
cloud-placeholder detection applies.

## No indexing/sync infrastructure

Live search at query time only: walk the allowlisted directories, read
matching files, do a plain case-insensitive substring match over
content. No persistent search index is built or maintained, matching
this whole package's established "no indexing pipeline in V1"
philosophy (see `mcp_connectors/__init__.py`). **Honesty note**: this
means an office/PDF file is re-parsed from scratch on every single query
that walks past it -- there is no cache. That's markedly slower than the
plain-text read path (a multi-page PDF or a large `.xlsx` is real parsing
work, not a `read_text()` call), and gets slower the more office/PDF
files live in an allowlisted directory. Fine for V1's "no indexing
pipeline" scope and for the directory sizes this is designed around, but
a real cost, not a free extension of the existing design -- worth an
indexing/caching pass if this connector is ever pointed at a directory
with hundreds of large office documents.

## V1 file-type scope: opt-in categories, not a fixed extension list

Every connector's config declares a `file_categories` list (default:
`["text"]` only) -- the config-time "permission" gate for *what kinds of
files* this connector may ever read, layered on top of (not a
replacement for) the directory allowlist described below. Four
categories exist, defined by `FILE_CATEGORY_EXTENSIONS` below:

  - **`"text"`** (the default, always safe, no library needed): `.md`,
    `.markdown`, `.txt`, `.rst`.
  - **`"code"`** (no library needed -- read via the exact same plain-text
    path as `"text"`): a reasonably broad, deliberately non-exhaustive
    set of common source/config file extensions -- see `CODE_EXTENSIONS`
    below for the full list (Python, JS/TS, Java, C/C++/C#, Go, Rust,
    Ruby, PHP, Swift, Kotlin, Scala, shell, SQL, YAML/JSON/TOML/XML,
    HTML/CSS, and a handful more).
  - **`"office"`** (needs the `documents` extra): `.docx`, `.xlsx`,
    `.pptx` -- parsed via `python-docx`/`openpyxl`/`python-pptx`
    respectively, extracting real embedded text (paragraphs + table
    cells for `.docx`; cell values across all sheets for `.xlsx`; text
    frame text across all slide shapes for `.pptx`).
  - **`"pdf"`** (needs the `documents` extra): `.pdf` -- parsed via
    `pypdf`, extracting each page's text and joining them.

  **Explicitly, deliberately out of scope: OCR / image-based text
  recognition.** This was a direct question put to the project owner --
  "does this include images via OCR?" -- answered explicitly:
  structured documents with real embedded text only, not image/screenshot
  recognition. No Tesseract or other OCR dependency is used or planned
  here; a scanned/image-only PDF, or an image file of any kind (`.png`/
  `.jpg`/etc.), is never readable by this connector, in any category,
  confirmed as a deliberate exclusion rather than a gap.

  A category not listed in a connector's config is a hard `search()`/
  `fetch()` boundary, identically to the directory allowlist: a `.py`
  file sitting in an allowlisted directory is invisible to `search()`
  and refused by `fetch()` with a clear error unless `"code"` is in that
  connector's configured `file_categories` -- see
  `local_docs/config.py`/`onedrive/config.py`'s `file_categories` field
  and `_validate_file_categories`. This is backward compatible on
  purpose: an existing config file with no `file_categories` key at all
  gets the pydantic field default (`["text"]`) and behaves exactly as it
  did before office/PDF/code support existed -- widening scope always
  requires an explicit config edit, never happens automatically from a
  package upgrade.

  Requesting `"office"` or `"pdf"` when the corresponding library isn't
  importable in the current environment is a **config-validation-time**
  error (`missing_libraries_for_categories`, called from each
  connector's config model), not a confusing failure discovered later on
  the first matching file -- see `pyproject.toml`'s `documents` extra.
  `"office"` requires all three of `python-docx`/`openpyxl`/
  `python-pptx` to be importable (not just whichever one a particular
  file happens to need) -- deliberately simple/predictable rather than a
  category that silently half-works depending on which of the three
  libraries happens to be installed.

## The precision requirement, adapted for local files

Every other connector in this package (Jira/Confluence/SharePoint)
enforces scope in two places: a config-time hard allowlist, and a
query-time restriction expressed in the backend's own native query
language (JQL/CQL/Graph Search `Path:` clauses) -- see
`mcp_connectors/__init__.py`. There is no remote API doing server-side
scoping for local files, so this module is fully responsible for the
equivalent guarantee, in two analogous places:

  1. **Config-time hard allowlist of directories** -- each connector's
     config model (`local_docs/config.py::LocalDocsConnectorConfig`,
     `onedrive/config.py::OneDriveConnectorConfig`) resolves every
     configured directory via `Path.resolve(strict=True)` at
     config-*validation* time (a `pydantic` `field_validator`, not
     something this module does itself) -- so `config.allowed_directories`
     already holds canonical, existing, real absolute paths by the time
     any client/search code sees them. `enforce_allowlist` (from
     `mcp_connectors/common.py`, reused as-is -- not reimplemented) is
     then used at query time by each connector's `client.py` exactly the
     way Jira/Confluence use it for project/space keys, to validate any
     caller-supplied `directories` subset against that config-time list.

  2. **Query-time path-safety enforcement -- new logic, specific to this
     connector family, written and tested here for the first time.**
     A directory-prefix string check alone is not a real security
     boundary: a symlink *inside* an allowed directory can point
     *outside* it, and a `fetch()` id is caller-supplied, opaque input
     that could contain a path-traversal sequence (e.g. `"../../etc/
     passwd"`). So before any file's content is read or returned in a
     result, **`resolve_within_allowlist`/`require_within_allowlist`
     below always resolve the file's *real* path** (`Path.resolve()`,
     which follows symlinks) and verify that resolved path is actually
     `relative_to` one of the allowlisted (already-resolved) directories
     -- never a bare `str.startswith(...)` prefix check, which a symlink
     or a `..`-laden path can defeat. A path that resolves outside every
     allowed directory -- whether via a traversal sequence in a `fetch()`
     id, or via a symlink encountered mid-walk during `search()` -- is
     never read and never appears in a result. `search()`'s walk drops
     such entries silently (a bad symlink somewhere in a big directory
     tree shouldn't abort an otherwise-valid search); `fetch()`'s
     `require_within_allowlist` raises `ConnectorAPIError` instead,
     because a `fetch()` id is a specific, caller-supplied request that
     deserves a loud, explicit refusal rather than a silent empty result.
     See `tests/test_local_fs.py` for the symlink-escape and
     path-traversal tests this guarantee is checked against directly --
     not just asserted in prose here. **This check is fully agnostic to
     file type/category** -- it runs identically, before any
     format-specific extraction function below is ever called, whether
     the target is a `.txt` file or a `.pdf`; a symlink escape pointing
     at a `.pdf` is rejected by the exact same `require_within_allowlist`
     call as one pointing at a `.txt` file (see
     `test_symlink_escape_is_rejected_for_office_and_pdf_files`).

## Per-format text extraction: defensive against corrupted/malformed files

`_extract_docx_text`/`_extract_xlsx_text`/`_extract_pptx_text`/
`_extract_pdf_text` each wrap their underlying library's parse call in a
broad `except Exception` -- deliberately broad, not a specific exception
class, because each library raises a *different* exception type for
malformed/corrupted/wrong-format input (verified directly against the
installed libraries, not assumed: `pypdf.errors.PdfStreamError` for a
truncated/garbage PDF, `docx.opc.exceptions.PackageNotFoundError` for a
non-`.docx`/corrupted file, `zipfile.BadZipFile` for a corrupted
`.xlsx`, `pptx.exc.PackageNotFoundError` for a corrupted `.pptx` -- and
that's not an exhaustive list of what a truly adversarial or just
bit-rotted file could trigger). Any of them means "this one file can't
be read," not "abort the whole search" -- exactly the same posture
`_read_text` already has for a plain-text file it can't decode/open
(`OSError` -> `None`). A password-protected PDF is treated the same way
(`reader.is_encrypted` -> skip; no password-guessing attempted). None of
this is speculative -- see `tests/test_local_fs.py`'s
`test_*_corrupted_file_is_skipped_not_crashed` tests, which feed each
extractor a real file with the right extension but garbage bytes as
content and assert `None` comes back, not a raised exception.

## OneDrive Files-On-Demand: a real, honest gap, not a false "fully handled"

OneDrive's "Files On-Demand" feature can leave a file as a cloud-only
placeholder -- present in a directory listing, but with no actual local
content until it's opened/downloaded through the OneDrive client or
Explorer/Finder UI. Returning such a file's "content" naively would
either read back an empty string (silently wrong -- looks like a real
empty document) or, in principle, whatever partial/placeholder bytes the
OS happens to expose. `looks_like_cloud_only_placeholder` below is a
**best-effort, explicitly partial** detector, applied only when a
connector opts in (`detect_cloud_placeholders=True` -- `onedrive`'s
client passes this; `local_docs`'s does not, since a genuinely empty
`.txt`/`.md` file in a plain local directory is a normal thing to have
and shouldn't be treated as suspicious):

  - **Windows**: `os.stat()` exposes `st_file_attributes` on Windows
    only. This checks it for `FILE_ATTRIBUTE_RECALL_ON_DATA_ACCESS`
    (`0x00400000`) or `FILE_ATTRIBUTE_OFFLINE` (`0x00001000`) -- real,
    documented Win32 file-attribute bits OneDrive sets on a cloud-only
    placeholder that hasn't been hydrated to local disk. This is the one
    signal here with real platform-attribute backing, not just a
    heuristic.
  - **All platforms, as a second-line heuristic**: a zero-byte file is
    treated as a suspected placeholder. A genuinely empty file is
    possible but unusual for a synced document, so this errs toward
    flagging it explicitly rather than silently returning empty content
    as if it were a confirmed-real empty document.

  **What this does *not* attempt** (a real gap, documented here rather
  than glossed over): on **macOS**, OneDrive's Files-On-Demand is
  implemented via Apple's File Provider extension, and a file's real
  download status lives behind
  `NSMetadataItemUbiquitousItemDownloadingStatusKey` / the File Provider
  APIs -- reachable from Cocoa/Foundation, **not** from anything in
  Python's stdlib `os`/`pathlib`. Detecting it properly would need a
  native bridge (e.g. `pyobjc`), which this connector deliberately does
  not take on as a dependency for a stdlib-only V1. A macOS cloud-only
  placeholder that happens to be non-zero-byte (uncommon, but not
  impossible) will **not** be caught by this module and could still be
  read as if it were real content. On Linux, OneDrive has no first-party
  client, so this gap doesn't really apply there, but the zero-byte
  heuristic still runs for whatever's on disk.

  Where this check fires: `search_local_files` silently skips a detected
  placeholder (it contributes no real searchable text anyway); `fetch()`
  raises a clear `ConnectorAPIError` naming the file and explaining it
  looks like an undownloaded placeholder, rather than returning empty or
  partial content as if it were the real thing.
"""
from __future__ import annotations

import os
from pathlib import Path
from typing import Dict, FrozenSet, Iterator, List, Optional, Sequence, Tuple

from mcp_connectors.common import ConnectorAPIError, Document

# -- deferred imports for the office/pdf parsing libraries ---------------------
#
# Mirrors `mcp_connectors/common.py`'s `keyring` deferred-ImportError
# pattern: this module must import cleanly even when the `documents`
# extra isn't installed (e.g. a `local_docs`/`onedrive` install that only
# ever wants `"text"`/`"code"` categories, which need no third-party
# library at all). Each library's ImportError, if any, is captured once
# at module-import time and only ever turned into a real, user-facing
# error at config-validation time (`missing_libraries_for_categories`,
# called from each connector's `file_categories` field validator) or
# defensively inside the corresponding `_extract_*_text` function.

try:
    import docx as _docx  # python-docx

    _DOCX_IMPORT_ERROR: Optional[BaseException] = None
except ImportError as exc:  # pragma: no cover - exercised only without `documents`
    _docx = None
    _DOCX_IMPORT_ERROR = exc

try:
    import openpyxl as _openpyxl

    _OPENPYXL_IMPORT_ERROR: Optional[BaseException] = None
except ImportError as exc:  # pragma: no cover - exercised only without `documents`
    _openpyxl = None
    _OPENPYXL_IMPORT_ERROR = exc

try:
    import pptx as _pptx  # python-pptx

    _PPTX_IMPORT_ERROR: Optional[BaseException] = None
except ImportError as exc:  # pragma: no cover - exercised only without `documents`
    _pptx = None
    _PPTX_IMPORT_ERROR = exc

try:
    import pypdf as _pypdf

    _PYPDF_IMPORT_ERROR: Optional[BaseException] = None
except ImportError as exc:  # pragma: no cover - exercised only without `documents`
    _pypdf = None
    _PYPDF_IMPORT_ERROR = exc


# -- file-type scope: categories, extensions -------------------------------------

#: `"text"` category -- see module docstring. Matched case-insensitively
#: against a candidate file's suffix. Kept under its original name for
#: backward compatibility (this constant existed before `file_categories`
#: did, and is still exactly the default-category extension set).
PLAIN_TEXT_EXTENSIONS = frozenset({".md", ".markdown", ".txt", ".rst"})
TEXT_EXTENSIONS = PLAIN_TEXT_EXTENSIONS

#: `"code"` category -- see module docstring. A deliberately non-exhaustive
#: but reasonably broad set of common source/config file extensions,
#: read via the exact same plain-text path as `TEXT_EXTENSIONS` (no
#: library needed for this category).
CODE_EXTENSIONS = frozenset(
    {
        ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".c", ".h", ".cc", ".cpp", ".hpp",
        ".cs", ".go", ".rs", ".rb", ".php", ".swift", ".kt", ".kts", ".scala",
        ".sh", ".bash", ".zsh", ".sql",
        ".yaml", ".yml", ".json", ".toml", ".xml", ".ini", ".cfg", ".conf",
        ".html", ".htm", ".css", ".scss", ".less",
        ".lua", ".pl", ".pm", ".r", ".dart", ".ex", ".exs", ".clj", ".cljs", ".hs",
    }
)

#: `"office"` category -- see module docstring. Each extension maps to
#: exactly one parsing library (`.docx` -> `python-docx`, `.xlsx` ->
#: `openpyxl`, `.pptx` -> `python-pptx`), all three required to be
#: importable for the category as a whole -- see
#: `missing_libraries_for_categories`.
OFFICE_EXTENSIONS = frozenset({".docx", ".xlsx", ".pptx"})

#: `"pdf"` category -- see module docstring.
PDF_EXTENSIONS = frozenset({".pdf"})

#: The four known `file_categories` values -- also the `Literal` choices
#: on each connector's config model (`local_docs/config.py`,
#: `onedrive/config.py`), imported from here so both stay in sync with
#: this module rather than redeclaring the list.
FileCategory = str  # `Literal["text", "code", "office", "pdf"]` at the config layer
KNOWN_FILE_CATEGORIES: Tuple[str, ...] = ("text", "code", "office", "pdf")

#: The config-model field default -- unchanged, existing-file-scope
#: behavior for any config that predates `file_categories` entirely (see
#: module docstring's backward-compatibility note).
DEFAULT_FILE_CATEGORIES: Tuple[str, ...] = ("text",)

FILE_CATEGORY_EXTENSIONS: Dict[str, FrozenSet[str]] = {
    "text": TEXT_EXTENSIONS,
    "code": CODE_EXTENSIONS,
    "office": OFFICE_EXTENSIONS,
    "pdf": PDF_EXTENSIONS,
}

#: Win32 `FILE_ATTRIBUTE_*` bits, only ever present on `os.stat_result
#: .st_file_attributes` on Windows -- see module docstring's OneDrive
#: section. Named here as plain module-level ints (not imported from
#: anywhere platform-specific) so this module still imports cleanly on
#: macOS/Linux; the attribute simply won't be present on those
#: platforms' `stat` results, handled via `getattr(..., None)` below.
_FILE_ATTRIBUTE_RECALL_ON_DATA_ACCESS = 0x00400000
_FILE_ATTRIBUTE_OFFLINE = 0x00001000


def extensions_for_categories(categories: Sequence[str]) -> FrozenSet[str]:
    """The file-extension allowlist a given `file_categories` config
    resolves to -- e.g. `["text", "code"]` -> the union of
    `TEXT_EXTENSIONS` and `CODE_EXTENSIONS`. An unknown category name is
    silently ignored here (pydantic's `Literal` field type is what
    actually rejects an unknown category value at config-validation
    time -- this function is also called from test code with plain
    strings, so it stays permissive rather than raising)."""
    result: set = set()
    for category in categories:
        result |= FILE_CATEGORY_EXTENSIONS.get(category, frozenset())
    return frozenset(result)


def missing_libraries_for_categories(categories: Sequence[str]) -> List[str]:
    """Returns the human-readable `pip`-package names of any parsing
    library required by `categories` that isn't actually importable in
    this environment right now -- e.g. `["office"]` on an environment
    without the `documents` extra installed returns
    `["python-docx", "openpyxl", "python-pptx"]`. Empty list means every
    requested category's dependencies (if any) are satisfied (`"text"`/
    `"code"` never contribute anything here -- they need nothing beyond
    stdlib `pathlib`). Called from each connector's `file_categories`
    config-field validator, so a config requesting `"office"`/`"pdf"`
    support fails fast and clearly at config-load time rather than
    silently skipping every matching file at query time."""
    missing: List[str] = []
    if "office" in categories:
        if _DOCX_IMPORT_ERROR is not None:
            missing.append("python-docx")
        if _OPENPYXL_IMPORT_ERROR is not None:
            missing.append("openpyxl")
        if _PPTX_IMPORT_ERROR is not None:
            missing.append("python-pptx")
    if "pdf" in categories:
        if _PYPDF_IMPORT_ERROR is not None:
            missing.append("pypdf")
    return missing


def _real_path_within_allowlist(path: Path, allowed_dirs: Sequence[Path]) -> Optional[Tuple[Path, Path]]:
    """The core path-safety check -- see module docstring's "precision
    requirement" section. Resolves `path`'s real, symlink-followed
    location and checks it's actually inside one of `allowed_dirs`
    (themselves assumed already-resolved, real paths -- see each
    connector's config model). Returns `(real_path, containing_allowed_dir)`
    if safe, `None` if `path` doesn't resolve to anything inside any
    allowed directory (including: doesn't exist, isn't accessible, or
    resolves -- directly or via a symlink -- outside every allowed dir).
    Never raises; callers decide whether a `None` result means "skip
    silently" (`search`'s directory walk) or "refuse loudly"
    (`require_within_allowlist`, used by `fetch()`). This check has no
    awareness of file type/category whatsoever -- it runs identically
    regardless of what's ultimately going to try to read the file
    afterward."""
    try:
        real = path.resolve(strict=True)
    except OSError:
        return None
    for allowed in allowed_dirs:
        try:
            real.relative_to(allowed)
        except ValueError:
            continue
        return real, allowed
    return None


def require_within_allowlist(path: Path, allowed_dirs: Sequence[Path]) -> Tuple[Path, Path]:
    """The loud-refusal counterpart to `_real_path_within_allowlist`,
    used by `fetch_local_file` for caller-supplied ids. Raises
    `ConnectorAPIError` -- never silently resolves and reads -- if `path`
    doesn't resolve to a real location inside one of `allowed_dirs`. This
    is what makes a path-traversal id (e.g. `"../../etc/passwd"`) or a
    symlink pointing outside the allowlist a hard, explicit failure
    rather than a quietly-narrowed no-op."""
    result = _real_path_within_allowlist(path, allowed_dirs)
    if result is None:
        raise ConnectorAPIError(
            f"{str(path)!r} does not resolve to a real, existing path inside any "
            "of this connector's allowlisted directories -- refused. This "
            "connector rejects path-traversal-style ids and symlinks that "
            "resolve outside the configured allowlist rather than silently "
            "following them."
        )
    return result


def looks_like_cloud_only_placeholder(real_path: Path) -> bool:
    """Best-effort OneDrive Files-On-Demand cloud-only-placeholder
    detector -- see module docstring's dedicated section for exactly
    what this does and does not catch. `real_path` must already be a
    resolved, verified-safe path (i.e. the output of
    `_real_path_within_allowlist`/`require_within_allowlist`), not raw
    caller input."""
    try:
        st = real_path.stat()
    except OSError:
        return False

    file_attributes = getattr(st, "st_file_attributes", None)
    if file_attributes is not None and file_attributes & (
        _FILE_ATTRIBUTE_RECALL_ON_DATA_ACCESS | _FILE_ATTRIBUTE_OFFLINE
    ):
        return True

    return st.st_size == 0


def iter_candidate_files(
    allowed_dirs: Sequence[Path], allowed_extensions: FrozenSet[str] = PLAIN_TEXT_EXTENSIONS
) -> Iterator[Tuple[Path, Path]]:
    """Walk every allowlisted directory, yielding `(real_path,
    containing_allowed_dir)` for every regular file whose extension is
    in `allowed_extensions` (defaults to the `"text"` category's
    extensions, for backward compatibility with callers that predate
    `file_categories` -- in practice every real caller now passes
    `extensions_for_categories(config.file_categories)` explicitly) --
    symlink-escape-guarded (see module docstring): a symlink resolving
    outside every allowed directory is silently dropped, not yielded.
    Directories that vanish or become unreadable between config-load
    time and a query (e.g. an unmounted removable drive) are skipped,
    not a hard failure of the whole search -- a `search()` call
    intentionally degrades gracefully to "found nothing there" rather
    than aborting the entire multi-directory search over one bad entry.
    Sorted for deterministic ordering (tests, and stable "first N
    results" behavior under `result_limit`)."""
    for directory in allowed_dirs:
        try:
            if not directory.is_dir():
                continue
            candidates = sorted(directory.rglob("*"))
        except OSError:
            continue
        for candidate in candidates:
            try:
                if not candidate.is_file():
                    continue
            except OSError:
                continue
            if candidate.suffix.lower() not in allowed_extensions:
                continue
            resolved = _real_path_within_allowlist(candidate, allowed_dirs)
            if resolved is None:
                continue
            yield resolved


def _read_text(real_path: Path) -> Optional[str]:
    try:
        return real_path.read_text(encoding="utf-8", errors="replace")
    except OSError:
        return None


def _extract_pdf_text(real_path: Path) -> Optional[str]:
    """Extracts page text from a `.pdf` via `pypdf`. Returns `None`
    (never raises) for: the `documents` extra not being installed, an
    encrypted/password-protected PDF (no password-guessing attempted --
    this connector has no credential of any kind, on principle), or a
    corrupted/malformed/wrong-format file -- see module docstring's
    "Per-format text extraction" section for exactly which real
    exception types were observed and why this catches broadly rather
    than one specific class."""
    if _PYPDF_IMPORT_ERROR is not None:
        return None
    try:
        reader = _pypdf.PdfReader(str(real_path))
        if getattr(reader, "is_encrypted", False):
            return None
        parts: List[str] = []
        for page in reader.pages:
            text = page.extract_text() or ""
            if text:
                parts.append(text)
        return "\n".join(parts)
    except Exception:  # noqa: BLE001 - pypdf raises several different exception
        # types for malformed/truncated/non-PDF bytes (observed directly:
        # PdfStreamError, PdfReadError, plus lower-level struct/zlib errors
        # depending on exactly what's corrupted) -- any of them means "can't
        # read this one file," not "abort the whole search." See module
        # docstring.
        return None


def _extract_docx_text(real_path: Path) -> Optional[str]:
    """Extracts paragraph and table-cell text from a `.docx` via
    `python-docx`. See `_extract_pdf_text`'s docstring for the general
    "returns `None`, never raises" contract this and the other
    `_extract_*_text` functions all share."""
    if _DOCX_IMPORT_ERROR is not None:
        return None
    try:
        document = _docx.Document(str(real_path))
        parts: List[str] = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        parts.append(cell.text)
        return "\n".join(parts)
    except Exception:  # noqa: BLE001 - python-docx raises
        # docx.opc.exceptions.PackageNotFoundError for a corrupted/non-.docx
        # file (observed directly), among other exceptions a malformed zip/XML
        # payload could trigger. See module docstring.
        return None


def _extract_xlsx_text(real_path: Path) -> Optional[str]:
    """Extracts every non-empty cell value across every worksheet of a
    `.xlsx` via `openpyxl`, in `read_only`/`data_only` mode (reads
    computed values, not formula source text). See `_extract_pdf_text`'s
    docstring for the shared "returns `None`, never raises" contract."""
    if _OPENPYXL_IMPORT_ERROR is not None:
        return None
    workbook = None
    try:
        workbook = _openpyxl.load_workbook(str(real_path), read_only=True, data_only=True)
        parts: List[str] = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for value in row:
                    if value is not None and str(value).strip():
                        parts.append(str(value))
        return "\n".join(parts)
    except Exception:  # noqa: BLE001 - openpyxl raises zipfile.BadZipFile for a
        # corrupted file and openpyxl.utils.exceptions.InvalidFileException for
        # a wrong-format one (both observed directly), among others. See
        # module docstring.
        return None
    finally:
        if workbook is not None:
            try:
                workbook.close()
            except Exception:  # noqa: BLE001 - closing a partially-loaded/
                # already-broken workbook must never itself raise past this
                # defensive extraction function.
                pass


def _extract_pptx_text(real_path: Path) -> Optional[str]:
    """Extracts text-frame text across every shape of every slide of a
    `.pptx` via `python-pptx`. See `_extract_pdf_text`'s docstring for
    the shared "returns `None`, never raises" contract."""
    if _PPTX_IMPORT_ERROR is not None:
        return None
    try:
        presentation = _pptx.Presentation(str(real_path))
        parts: List[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs)
                    if text:
                        parts.append(text)
        return "\n".join(parts)
    except Exception:  # noqa: BLE001 - python-pptx raises
        # pptx.exc.PackageNotFoundError for a corrupted/non-.pptx file
        # (observed directly), among others. See module docstring.
        return None


def _extract_text(real_path: Path) -> Optional[str]:
    """Dispatches to the right extraction path by extension -- plain
    `_read_text` for `"text"`/`"code"`-category files, the matching
    `_extract_*_text` function for an office/PDF file. Only ever called
    on an already-`require_within_allowlist`/`_real_path_within_allowlist`
    -verified real path (see module docstring's "precision requirement"
    section) -- this function itself does no path-safety checking, by
    design, since that check is fully type-agnostic and already done by
    the time any caller reaches here."""
    suffix = real_path.suffix.lower()
    if suffix in TEXT_EXTENSIONS or suffix in CODE_EXTENSIONS:
        return _read_text(real_path)
    if suffix == ".pdf":
        return _extract_pdf_text(real_path)
    if suffix == ".docx":
        return _extract_docx_text(real_path)
    if suffix == ".xlsx":
        return _extract_xlsx_text(real_path)
    if suffix == ".pptx":
        return _extract_pptx_text(real_path)
    return None  # pragma: no cover - unreachable given upstream extension filtering


def _to_document(*, real_path: Path, container: Path, text: str, source: str, query: Optional[str]) -> Document:
    last_modified = None
    try:
        last_modified = _mtime_to_datetime(real_path.stat().st_mtime)
    except OSError:
        pass

    snippet = _build_snippet(text, query)
    return Document(
        id=str(real_path),
        title=real_path.name,
        snippet=snippet,
        source=source,
        url=real_path.as_uri(),
        last_modified=last_modified,
        container=str(container),
        metadata={"extension": real_path.suffix.lower(), "size_bytes": len(text.encode("utf-8"))},
    )


def _mtime_to_datetime(mtime: float):
    from datetime import datetime, timezone

    return datetime.fromtimestamp(mtime, tz=timezone.utc)


def _build_snippet(text: str, query: Optional[str], *, context_chars: int = 200, max_len: int = 2000) -> str:
    """A short, human-useful snippet: text around the first match of
    `query` if given (case-insensitive), else just the start of the
    file. Always capped at `max_len` -- mirrors the other connectors'
    "cap snippet length" convention (e.g. Jira's `snippet[:2000]`)."""
    if query:
        lower_text = text.lower()
        idx = lower_text.find(query.lower())
        if idx != -1:
            start = max(0, idx - context_chars)
            end = min(len(text), idx + len(query) + context_chars)
            prefix = "..." if start > 0 else ""
            suffix = "..." if end < len(text) else ""
            return (prefix + text[start:end] + suffix)[:max_len]
    return text[:max_len]


def search_local_files(
    query: str,
    allowed_dirs: Sequence[Path],
    *,
    source: str,
    limit: int,
    detect_cloud_placeholders: bool = False,
    file_categories: Sequence[str] = DEFAULT_FILE_CATEGORIES,
) -> List[Document]:
    """Live, no-index text search over `allowed_dirs` (already
    config-time-resolved and, if this is a narrowed request, already
    `enforce_allowlist`-validated by the caller -- this function itself
    trusts `allowed_dirs` as the scope to search, exactly like
    `JiraClient.search()` trusts the JQL it's handed). Case-insensitive
    plain substring match over each candidate file's extracted content
    (plain read for `"text"`/`"code"`, real parsing for `"office"`/
    `"pdf"` -- see `_extract_text`). Capped at `limit` results, stopping
    the walk early once reached (never an unbounded scan-everything-
    then-cap). `file_categories` defaults to `("text",)` -- the same
    default as each connector's config model -- for any caller that
    predates `file_categories` existing at all."""
    query_norm = (query or "").strip()
    if not query_norm:
        raise ConnectorAPIError("query text must not be empty")

    allowed_extensions = extensions_for_categories(file_categories)
    needle = query_norm.lower()
    results: List[Document] = []
    for real_path, container in iter_candidate_files(allowed_dirs, allowed_extensions):
        if detect_cloud_placeholders and looks_like_cloud_only_placeholder(real_path):
            # A cloud-only placeholder has no real local content to
            # search -- skip it gracefully rather than matching on
            # nothing or raising mid-walk. See module docstring.
            continue
        text = _extract_text(real_path)
        if text is None:
            # Either a plain-text file this process can't read, or an
            # office/PDF file that failed to parse (corrupted, encrypted,
            # or the `documents` extra isn't installed) -- either way,
            # skip this one file, don't abort the walk. See module
            # docstring's "Per-format text extraction" section.
            continue
        if needle not in text.lower():
            continue
        results.append(_to_document(real_path=real_path, container=container, text=text, source=source, query=query_norm))
        if len(results) >= limit:
            break
    return results


def fetch_local_file(
    file_id: str,
    allowed_dirs: Sequence[Path],
    *,
    source: str,
    detect_cloud_placeholders: bool = False,
    file_categories: Sequence[str] = DEFAULT_FILE_CATEGORIES,
) -> Document:
    """Fetch one file by id (an absolute path string, normally a
    previous `search_local_files` result's `Document.id`). Path-safety
    is enforced via `require_within_allowlist` -- see module docstring --
    which is what makes a path-traversal or symlink-escape id a hard,
    explicit `ConnectorAPIError` rather than a silently-resolved read,
    *before* any type/category check or extraction is even attempted.
    `file_categories` defaults to `("text",)`, matching
    `search_local_files`."""
    raw = (file_id or "").strip()
    if not raw:
        raise ConnectorAPIError("file id must not be empty")

    real_path, container = require_within_allowlist(Path(raw).expanduser(), allowed_dirs)

    if not real_path.is_file():
        raise ConnectorAPIError(f"{file_id!r} does not resolve to a regular file")

    allowed_extensions = extensions_for_categories(file_categories)
    if real_path.suffix.lower() not in allowed_extensions:
        raise ConnectorAPIError(
            f"{file_id!r} has extension {real_path.suffix!r}, which is outside this "
            f"connector's currently-enabled file_categories ({list(file_categories)!r}, "
            f"resolving to extensions {sorted(allowed_extensions)!r}). Add the right "
            "category ('code'/'office'/'pdf') to this connector's config to grant "
            "access -- see local_fs/search.py's module docstring for what each "
            "category covers. OCR/image-based text extraction is out of scope "
            "regardless of configured categories -- a deliberate exclusion, not a gap."
        )
    if detect_cloud_placeholders and looks_like_cloud_only_placeholder(real_path):
        raise ConnectorAPIError(
            f"{file_id!r} looks like a OneDrive Files-On-Demand cloud-only "
            "placeholder with no local content yet (zero-byte, or flagged "
            "offline/recall-on-access on Windows) -- open it via the OneDrive "
            "client, Explorer, or Finder first to download it, then fetch again. "
            "This detection is best-effort -- see mcp_connectors/local_fs/search.py's "
            "module docstring for exactly what is and isn't checked."
        )

    text = _extract_text(real_path)
    if text is None:
        raise ConnectorAPIError(
            f"failed to read/parse {file_id!r} ({real_path}) -- the file may be "
            "corrupted, password-protected/encrypted, in an unexpected format for "
            "its extension, or (for 'office'/'pdf' categories) the `documents` "
            "extra may not be installed in this environment."
        )
    return _to_document(real_path=real_path, container=container, text=text, source=source, query=None)
